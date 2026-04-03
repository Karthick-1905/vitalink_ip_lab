import sys
from pathlib import Path

import pandas as pd

EXTRA_SITE_PACKAGES = Path("/tmp/codex_pptx")
if EXTRA_SITE_PACKAGES.exists():
    sys.path.insert(0, str(EXTRA_SITE_PACKAGES))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "output"
DOCS_DIR = BASE_DIR / "docs"

PPTX_PATH = DOCS_DIR / "warfarin_program_slides.pptx"

DOSE_METRICS = pd.read_csv(OUTPUT_DIR / "iwpc_model_metrics.csv")
DOSE_SHAP = pd.read_csv(OUTPUT_DIR / "iwpc_shap_top_features.csv")
STABILITY_METRICS = pd.read_csv(OUTPUT_DIR / "stability_model_metrics.csv")
STABILITY_SHAP = pd.read_csv(OUTPUT_DIR / "stability_shap_top_features.csv")
UPDATE_METRICS = pd.read_csv(OUTPUT_DIR / "continual_update_metrics.csv")

BG = RGBColor(243, 236, 223)
INK = RGBColor(31, 27, 22)
RUST = RGBColor(155, 77, 46)
OLIVE = RGBColor(63, 90, 78)
GOLD = RGBColor(179, 138, 60)
SLATE = RGBColor(92, 86, 79)
WHITE = RGBColor(255, 251, 245)
SOFT = RGBColor(231, 221, 205)
PALE = RGBColor(248, 243, 236)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.15), Inches(0.15), Inches(13.03), Inches(7.2))
    frame.fill.background()
    frame.line.color.rgb = SOFT
    frame.line.width = Pt(1.0)


def add_text(slide, left, top, width, height, text, size=20, color=INK, bold=False, font="Georgia", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    p.alignment = align
    return box


def add_paragraph(frame, text, size=16, color=INK, bold=False, font="Georgia", space_after=6):
    p = frame.add_paragraph()
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return p


def add_panel(slide, left, top, width, height, fill_rgb=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = SOFT
    shape.line.width = Pt(1)
    return shape


def add_header(slide, eyebrow, title):
    add_text(slide, Inches(0.5), Inches(0.38), Inches(3.5), Inches(0.22), eyebrow.upper(), 10, RUST, False, "Calibri")
    add_text(slide, Inches(0.5), Inches(0.65), Inches(11.9), Inches(0.72), title, 24, INK, True)


def add_footer(slide, label, number):
    add_text(slide, Inches(0.52), Inches(7.0), Inches(6.0), Inches(0.2), label, 9, SLATE, False, "Calibri")
    add_text(slide, Inches(12.2), Inches(7.0), Inches(0.4), Inches(0.2), f"{number:02d}", 9, SLATE, False, "Calibri", PP_ALIGN.RIGHT)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    dose_best = DOSE_METRICS.iloc[0]
    stability_best = STABILITY_METRICS.iloc[0]
    update_initial = UPDATE_METRICS.iloc[0]
    update_final = UPDATE_METRICS.iloc[-1]

    add_header(slide, "Warfarin Modeling Program", "Dose prediction, time-to-stability prediction, and continual model updates")
    lead_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(7.2), Inches(0.9))
    lead = lead_box.text_frame
    lead.word_wrap = True
    lead.text = ""
    add_paragraph(lead, "This deck consolidates the full workflow: benchmarking dose prediction against the IWPC calculator, evaluating time-to-stability models, and testing repeated retraining on shifted incoming cohorts.", 16, SLATE, False, space_after=0)

    cards = [
        ("01", "Dosage Prediction", "Best model slightly outperforms the IWPC pharmacogenetic calculator on RMSE."),
        ("02", "Time to Stability", "Current synthetic target remains nearly linear, so simple models hold up well."),
        ("03", "Continual Updates", "Retraining on shifted batches recovers large amounts of lost accuracy."),
    ]
    x_positions = [0.55, 4.22, 7.89]
    for (idx, title, body), x in zip(cards, x_positions):
        add_panel(slide, Inches(x), Inches(2.45), Inches(3.2), Inches(2.25), fill_rgb=PALE)
        add_text(slide, Inches(x + 0.18), Inches(2.68), Inches(0.5), Inches(0.2), idx, 10, RUST, False, "Calibri")
        add_text(slide, Inches(x + 0.18), Inches(2.95), Inches(2.75), Inches(0.4), title, 16, INK, True)
        body_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(3.4), Inches(2.72), Inches(0.95))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.text = ""
        add_paragraph(body_tf, body, 12, SLATE, False, font="Calibri", space_after=0)

    metric = add_panel(slide, Inches(10.0), Inches(1.35), Inches(2.75), Inches(3.35), fill_rgb=WHITE)
    metric.line.color.rgb = GOLD
    metric.line.width = Pt(1.5)
    add_text(slide, Inches(10.24), Inches(1.62), Inches(2.2), Inches(0.2), "BEST DOSE MODEL", 10, RUST, False, "Calibri")
    add_text(slide, Inches(10.24), Inches(1.9), Inches(2.2), Inches(0.55), str(dose_best["model"]), 22, OLIVE, True)
    add_text(slide, Inches(10.24), Inches(2.6), Inches(2.2), Inches(0.25), f"RMSE {dose_best['rmse']:.2f}", 15, INK, True, "Calibri")
    add_text(slide, Inches(10.24), Inches(2.95), Inches(2.2), Inches(0.25), f"Stability RMSE {stability_best['rmse']:.2f}d", 15, INK, True, "Calibri")
    add_text(slide, Inches(10.24), Inches(3.3), Inches(2.2), Inches(0.5), f"Dose adaptation {update_initial['dose_rmse']:.1f} -> {update_final['dose_rmse']:.1f}", 13, OLIVE, True, "Calibri")

    add_footer(slide, "Combined summary for all three requested tasks", 1)


def slide_dose_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Dose Prediction", "Dose benchmarking against IWPC and baseline machine-learning models")
    dose_best = DOSE_METRICS.iloc[0]
    dose_iwpc = DOSE_METRICS.loc[DOSE_METRICS["model"] == "IWPC Pharmacogenetic Calculator"].iloc[0]

    add_panel(slide, Inches(0.55), Inches(1.4), Inches(3.15), Inches(2.0), fill_rgb=PALE)
    add_text(slide, Inches(0.78), Inches(1.65), Inches(2.7), Inches(0.2), "TOP FINDING", 10, RUST, False, "Calibri")
    add_text(slide, Inches(0.78), Inches(1.95), Inches(2.7), Inches(0.45), str(dose_best["model"]), 20, OLIVE, True)
    add_text(slide, Inches(0.78), Inches(2.52), Inches(2.7), Inches(0.55), f"RMSE {dose_best['rmse']:.2f}, MAE {dose_best['mae']:.2f}, R² {dose_best['r2']:.3f}", 13, SLATE, False, "Calibri")
    add_text(slide, Inches(0.78), Inches(2.96), Inches(2.7), Inches(0.3), f"IWPC PGx RMSE {dose_iwpc['rmse']:.2f}, MAE {dose_iwpc['mae']:.2f}", 12, SLATE, False, "Calibri")

    add_panel(slide, Inches(0.55), Inches(3.58), Inches(3.15), Inches(2.55), fill_rgb=WHITE)
    add_text(slide, Inches(0.78), Inches(3.82), Inches(2.7), Inches(0.2), "TOP SHAP FEATURES", 10, RUST, False, "Calibri")
    y = 4.12
    for row in DOSE_SHAP.head(6).itertuples(index=False):
        add_text(slide, Inches(0.78), Inches(y), Inches(2.05), Inches(0.2), str(row.feature).replace("_", " "), 11, INK, False, "Calibri")
        add_text(slide, Inches(2.78), Inches(y), Inches(0.55), Inches(0.2), f"{row.mean_abs_shap:.3f}", 11, OLIVE, True, "Calibri", PP_ALIGN.RIGHT)
        y += 0.30

    table = slide.shapes.add_table(len(DOSE_METRICS) + 1, 6, Inches(3.95), Inches(1.4), Inches(8.4), Inches(4.74)).table
    headers = ["rank", "model", "rmse", "mae", "r2", "within_20_pct"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = SOFT
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(10)
        para.runs[0].font.bold = True
        para.runs[0].font.name = "Calibri"
    for r, row in enumerate(DOSE_METRICS.itertuples(index=False), start=1):
        values = [row.rank, row.model, f"{row.rmse:.4f}", f"{row.mae:.4f}", f"{row.r2:.4f}", f"{row.within_20_pct:.4f}"]
        for c, value in enumerate(values):
            para = table.cell(r, c).text_frame.paragraphs[0]
            para.text = str(value)
            para.runs[0].font.size = Pt(9)
            para.runs[0].font.name = "Calibri"

    add_footer(slide, "Dose prediction on the IWPC cohort", 2)


def slide_dose_visuals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Dose Visuals", "Leaderboard and fit diagnostics stay tightly clustered at the top")
    add_panel(slide, Inches(0.55), Inches(1.35), Inches(6.0), Inches(5.45), fill_rgb=WHITE)
    add_panel(slide, Inches(6.8), Inches(1.35), Inches(5.65), Inches(5.45), fill_rgb=WHITE)
    slide.shapes.add_picture(str(OUTPUT_DIR / "iwpc_model_comparison.png"), Inches(0.78), Inches(1.62), width=Inches(5.55), height=Inches(4.95))
    slide.shapes.add_picture(str(OUTPUT_DIR / "iwpc_prediction_scatter.png"), Inches(7.03), Inches(1.62), width=Inches(5.2), height=Inches(4.95))
    add_footer(slide, "Dose comparison and patient-level fit", 3)


def slide_stability_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Time to Stability", "The synthetic stability target is comparatively linear")
    best = STABILITY_METRICS.iloc[0]

    add_panel(slide, Inches(0.55), Inches(1.4), Inches(3.15), Inches(2.0), fill_rgb=PALE)
    add_text(slide, Inches(0.78), Inches(1.65), Inches(2.7), Inches(0.2), "BEST MODEL", 10, RUST, False, "Calibri")
    add_text(slide, Inches(0.78), Inches(1.95), Inches(2.7), Inches(0.45), str(best["model"]), 20, OLIVE, True)
    add_text(slide, Inches(0.78), Inches(2.52), Inches(2.7), Inches(0.6), f"RMSE {best['rmse']:.2f} days, MAE {best['mae']:.2f}, within 7 days {best['within_7_days_pct']:.1f}%", 13, SLATE, False, "Calibri")

    add_panel(slide, Inches(0.55), Inches(3.58), Inches(3.15), Inches(2.55), fill_rgb=WHITE)
    add_text(slide, Inches(0.78), Inches(3.82), Inches(2.7), Inches(0.2), "TOP SHAP FEATURES", 10, RUST, False, "Calibri")
    y = 4.12
    for row in STABILITY_SHAP.head(6).itertuples(index=False):
        add_text(slide, Inches(0.78), Inches(y), Inches(2.05), Inches(0.2), str(row.feature).replace("_", " "), 11, INK, False, "Calibri")
        add_text(slide, Inches(2.78), Inches(y), Inches(0.55), Inches(0.2), f"{row.mean_abs_shap:.3f}", 11, OLIVE, True, "Calibri", PP_ALIGN.RIGHT)
        y += 0.30

    table = slide.shapes.add_table(len(STABILITY_METRICS) + 1, 6, Inches(3.95), Inches(1.4), Inches(8.4), Inches(4.74)).table
    headers = ["rank", "model", "rmse", "mae", "r2", "within_7_days_pct"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = SOFT
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(10)
        para.runs[0].font.bold = True
        para.runs[0].font.name = "Calibri"
    for r, row in enumerate(STABILITY_METRICS.itertuples(index=False), start=1):
        values = [row.rank, row.model, f"{row.rmse:.4f}", f"{row.mae:.4f}", f"{row.r2:.4f}", f"{row.within_7_days_pct:.4f}"]
        for c, value in enumerate(values):
            para = table.cell(r, c).text_frame.paragraphs[0]
            para.text = str(value)
            para.runs[0].font.size = Pt(9)
            para.runs[0].font.name = "Calibri"
    add_footer(slide, "Time-to-stability comparison on the synthetic cohort", 4)


def slide_stability_visuals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Stability Visuals", "Scatter fit and SHAP views support a simpler operational model")
    positions = [0.55, 4.45, 8.35]
    images = ["stability_model_comparison.png", "stability_prediction_scatter.png", "stability_shap_bar.png"]
    for x, img in zip(positions, images):
        add_panel(slide, Inches(x), Inches(1.45), Inches(3.6), Inches(5.2), fill_rgb=WHITE)
        slide.shapes.add_picture(str(OUTPUT_DIR / img), Inches(x + 0.18), Inches(1.7), width=Inches(3.24), height=Inches(4.72))
    add_footer(slide, "Stability leaderboard, fit, and explainability", 5)


def slide_update_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Continual Updates", "Sequential retraining on shifted incoming cohorts materially restores performance")
    initial = UPDATE_METRICS.iloc[0]
    final = UPDATE_METRICS.iloc[-1]

    add_panel(slide, Inches(0.55), Inches(1.42), Inches(3.3), Inches(4.95), fill_rgb=PALE)
    add_text(slide, Inches(0.78), Inches(1.68), Inches(2.8), Inches(0.2), "DOSE ADAPTATION", 10, RUST, False, "Calibri")
    add_text(slide, Inches(0.78), Inches(1.98), Inches(2.8), Inches(0.25), f"RMSE {initial['dose_rmse']:.2f} -> {final['dose_rmse']:.2f}", 16, OLIVE, True, "Calibri")
    add_text(slide, Inches(0.78), Inches(2.32), Inches(2.8), Inches(0.25), f"Within 20% {initial['dose_within_20_pct']:.1f}% -> {final['dose_within_20_pct']:.1f}%", 14, SLATE, False, "Calibri")
    add_text(slide, Inches(0.78), Inches(3.02), Inches(2.8), Inches(0.2), "STABILITY ADAPTATION", 10, RUST, False, "Calibri")
    add_text(slide, Inches(0.78), Inches(3.32), Inches(2.8), Inches(0.25), f"RMSE {initial['stability_rmse']:.2f} -> {final['stability_rmse']:.2f}", 16, OLIVE, True, "Calibri")
    add_text(slide, Inches(0.78), Inches(3.66), Inches(2.8), Inches(0.25), f"Within 7 days {initial['stability_within_7_days_pct']:.1f}% -> {final['stability_within_7_days_pct']:.1f}%", 14, SLATE, False, "Calibri")

    table = slide.shapes.add_table(len(UPDATE_METRICS) + 1, 6, Inches(4.1), Inches(1.42), Inches(8.25), Inches(4.95)).table
    headers = ["round", "train_size", "dose_rmse", "dose_within_20_pct", "stability_rmse", "stability_within_7_days_pct"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = SOFT
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(10)
        para.runs[0].font.bold = True
        para.runs[0].font.name = "Calibri"
    for r, row in enumerate(UPDATE_METRICS.itertuples(index=False), start=1):
        values = [row.round, row.train_size, f"{row.dose_rmse:.4f}", f"{row.dose_within_20_pct:.4f}", f"{row.stability_rmse:.4f}", f"{row.stability_within_7_days_pct:.4f}"]
        for c, value in enumerate(values):
            para = table.cell(r, c).text_frame.paragraphs[0]
            para.text = str(value)
            para.runs[0].font.size = Pt(9)
            para.runs[0].font.name = "Calibri"
    add_footer(slide, "Continual-update experiment on shifted incoming data", 6)


def slide_update_visuals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Update Visuals", "The first update batch delivers the largest recovery, with smaller gains after that")
    positions = [0.55, 4.45, 8.35]
    images = ["continual_update_performance.png", "continual_update_shift.png", "continual_update_scatter.png"]
    for x, img in zip(positions, images):
        add_panel(slide, Inches(x), Inches(1.45), Inches(3.6), Inches(5.2), fill_rgb=WHITE)
        slide.shapes.add_picture(str(OUTPUT_DIR / img), Inches(x + 0.18), Inches(1.7), width=Inches(3.24), height=Inches(4.72))
    add_footer(slide, "Performance recovery, cohort shift, and pre/post fit", 7)


def slide_recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Recommendations", "What to deploy, what to keep simple, and what to operationalize")
    cards = [
        ("Deploy for Dose", "Use the tuned LightGBM dose model when the objective is best holdout accuracy against IWPC-style baselines."),
        ("Keep Stability Simple", "Use the linear stability model unless a richer real-world target becomes available."),
        ("Operational Retraining", "Schedule cumulative retraining when incoming cohorts diverge from the original patient mix."),
    ]
    positions = [0.55, 4.22, 7.89]
    for (title, body), x in zip(cards, positions):
        add_panel(slide, Inches(x), Inches(1.6), Inches(3.2), Inches(2.15), fill_rgb=PALE)
        add_text(slide, Inches(x + 0.18), Inches(1.9), Inches(2.8), Inches(0.3), title, 16, OLIVE, True)
        box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(2.28), Inches(2.75), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = ""
        add_paragraph(tf, body, 12, SLATE, False, "Calibri", 0)

    artifact = add_panel(slide, Inches(0.55), Inches(4.1), Inches(10.54), Inches(1.55), fill_rgb=WHITE)
    artifact.line.color.rgb = GOLD
    artifacts = [
        ("Master report", "docs/warfarin_full_program_report.md"),
        ("HTML deck", "docs/warfarin_program_slides.html"),
        ("PDF deck", "docs/warfarin_program_slides.pdf"),
        ("PPT deck", "docs/warfarin_program_slides.pptx"),
    ]
    x = 0.8
    for title, path in artifacts:
        add_text(slide, Inches(x), Inches(4.38), Inches(2.2), Inches(0.2), title.upper(), 9, RUST, False, "Calibri")
        add_text(slide, Inches(x), Inches(4.68), Inches(2.25), Inches(0.45), path, 11, INK, True, "Calibri")
        x += 2.55
    add_footer(slide, "Final recommendations and artifact map", 8)


def main():
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_dose_summary(prs)
    slide_dose_visuals(prs)
    slide_stability_summary(prs)
    slide_stability_visuals(prs)
    slide_update_summary(prs)
    slide_update_visuals(prs)
    slide_recommendations(prs)
    prs.save(PPTX_PATH)
    print(f"Wrote {PPTX_PATH}")


if __name__ == "__main__":
    main()
